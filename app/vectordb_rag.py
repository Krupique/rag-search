import docx
import PyPDF2
import sys
from pptx import Presentation
from os import listdir
from os.path import isfile, join, isdir
from langchain_text_splitters import TokenTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
from langchain_qdrant import Qdrant

def list_files(dir):
    # Initialize a empty list to store the file path
    arquivo_list = []
    
    # Iterate all files and directories on specif directory
    for f in listdir(dir):
        
        # If it's a file, add to the list
        if isfile(join(dir, f)):
            arquivo_list.append(join(dir, f))
        
        # If it is a directory, call the function recursively and add the results to the list
        elif isdir(join(dir, f)):
            arquivo_list += list_files(join(dir, f))
    
    # Return the list of files
    return arquivo_list


# Define the function that loads the text from a Word file 
def load_docx_text(arquivoname):
    
    # Open the Word file
    doc = docx.Document(arquivoname)
    
    # Extract the text from each paragraph and add it to the list
    fullText = [para.text for para in doc.paragraphs]
    
    
    # Join all texts into a single string separated by line breaks
    return '\n'.join(fullText)


# Define the function that load the text from a PowerPoint file
def load_pptx_text(arquivoname):
    
    # Open the PowerPoint file
    prs = Presentation(arquivoname)
    
    # Initialize a empty list to store the texts
    fullText = []
    
    # Iterate all slides
    for slide in prs.slides:
        
        # Iterate all shapes on slides
        for shape in slide.shapes:
            
            # If the shape has the "text" attributes, add it to the list 
            if hasattr(shape, "text"):
                fullText.append(shape.text)
    
    # Join all texts into a single string separated by line breaks
    return '\n'.join(fullText)



# Defines the main function for indexing documents
def main_indexing(mypath):
    
    # Sets the name of the model to be used to create the embeddings
    model_name = "sentence-transformers/msmarco-bert-base-dot-v5"
    model_kwargs = {'device': 'cpu'} # Model settings
    encode_kwargs = {'normalize_embeddings': True} # Set encoding settings

    # Initialize the HuggingFace embeddings class
    hf = HuggingFaceEmbeddings(model_name = model_name,
                               model_kwargs = model_kwargs,
                               encode_kwargs = encode_kwargs)

    client = QdrantClient("http://localhost:6333") # Initialize the Qdrant client
    collection_name = "VectorDB" # Sets the name of the embeddings collection

    if client.collection_exists(collection_name): # If the collection already exists, delete it
        client.delete_collection(collection_name)

    
    client.create_collection(collection_name, # Creates a new collection with specified parameters
                             vectors_config = VectorParams(size = 768, distance = Distance.DOT))

    qdrant = Qdrant(client, collection_name, hf) # Initialize the Qdrant instance

    #Prints a message informing that document indexing is starting
    print("\nIndexing the documents...\n")

    files_list = list_files(mypath)
    
    for file in files_list:
        
        try:
            # Initialize a empty string to store the file content
            file_content = ""
            
            if file_content.endswith(".pdf"):
                print("Indexing: " + file_content)
                reader = PyPDF2.PdfReader(file_content)
                for page in reader.pages:
                    file_content += " " + page.extract_text()
            
            elif file.endswith(".txt"):
                print("Indexing: " + file)
                with open(file, 'r') as f:
                    file_content = f.read()
            
            elif file.endswith(".docx"):
                print("Indexing: " + file)
                file_content = load_docx_text(file)
            
            elif file.endswith(".pptx"):
                print("Indexing: " + file)
                file_content = load_pptx_text(file)
            
            else:
                # If the file isn't of a supported format, continue to the next one.
                continue

            # Inicializa o divisor de texto com tamanho de chunk e sobreposição especificados
            text_splitter = TokenTextSplitter(chunk_size = 500, chunk_overlap = 50)
            
            # Spliting the file content in text chunks
            textos = text_splitter.split_text(file_content)
            
            # Create metadata for each text chunk
            # It allows that the LLM quotes the reference
            metadata = [{"path": file} for _ in textos]
            
            # Add the texts and its metadatas to Qdrant
            qdrant.add_texts(textos, metadatas = metadata)

        except Exception as e:
            # If an error occurs, print an error message
            print(f"The process failed for the file {file}: {e}")

    # Prints message informing that indexing is complete
    print("\nIndexing Completed!\n")


# Checks if the script is being executed directly
if __name__ == "__main__":

    # Gets command line arguments
    arguments = sys.argv
    
    # Checks if a directory path was provided
    if len(arguments) > 1:
        main_indexing(arguments[1])
    else:
        # If not, print an error message.
        print("You need to provide a path to the folder with documents to index.")

